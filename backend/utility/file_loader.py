import os
from typing import List
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PDFPlumberLoader
from pptx import Presentation

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=100,
    length_function=len,
    is_separator_regex=False
)

class PptxLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load_and_split(self) -> List[Document]:
        print(f"📊 Loading PPTX file: {self.file_path}")
        chunks = []
        try:
            ppt = Presentation(self.file_path)
            for slide in ppt.slides:
                full_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
                chunk_text = "\n".join(full_text)
                metadata = {"source": self.file_path}
                chunks.append(Document(page_content=chunk_text, metadata=metadata))
        except Exception as e:
            print(f"❌ Error reading PPTX file: {e}")
        return chunks

def load_pdf(file_path: str) -> List[Document]:
    print(f"📄 Loading PDF file: {file_path}")
    try:
        loader = PDFPlumberLoader(file_path)
        docs = loader.load_and_split()
        if not docs:
            print(f"⚠️ No content extracted from PDF: {file_path}")
        else:
            total_length = sum(len(doc.page_content.strip()) for doc in docs)
            print(f"✅ Loaded {len(docs)} chunks, total length: {total_length}")
        return docs
    except Exception as e:
        print(f"❌ Failed to load PDF: {e}")
        raise

def load_pptx(file_path: str) -> List[Document]:
    loader = PptxLoader(file_path)
    return loader.load_and_split()

def load_txt(file_path: str) -> List[Document]:
    print(f"📜 Loading TXT file: {file_path}")
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        docs = text_splitter.create_documents([text])
        for doc in docs:
            doc.metadata = {"source": file_path}
        print(f"✅ Loaded {len(docs)} chunks from TXT")
        return docs
    except Exception as e:
        print(f"❌ Error loading TXT file: {e}")
        raise

def load_and_split_file(file_path: str) -> List[Document]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext == ".pptx":
        return load_pptx(file_path)
    elif ext == ".txt":
        return load_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
